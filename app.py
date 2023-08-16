import streamlit as st
import os
import numpy as np
import pandas as pd
import mysql.connector
import plotly.express as px
import plotly.graph_objects as go
import plotly.io as pio
from wordcloud import WordCloud
import matplotlib.pyplot as plt
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

def create_binaryrep_with_images(images_folder):
    output = BytesIO()
    prs = Presentation()
    image_files = [file for file in os.listdir(images_folder) if file.endswith(('.jpg', '.jpeg', '.png', '.gif'))]
    for image_file in image_files:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(6)
        slide.shapes.add_picture(os.path.join(images_folder, image_file), left, top, width=width, height=height)

    prs.save(output)
    output_data = output.getvalue()
    return output_data
 

def save_plotly_figure_as_png(figure, file_path):
    fig = go.Figure(figure)
    pio.write_image(fig, file_path, format="png")

def download_plotly_figures_as_png(figures_list, folder_path):
    for i, figure in enumerate(figures_list):
        file_path = os.path.join(folder_path, f"figure_{i + 1}.png")
        save_plotly_figure_as_png(figure, file_path)

def retrieve_data_from_connections(connection):
        consolidated_df = pd.DataFrame()
        try:
            c = mysql.connector.connect(
                host=connection['hostname'],
                user=connection['username'],
                password=connection['password'],
                database=connection['database']
            )
            cursor = c.cursor()
            query = f"SELECT * FROM {connection['table']}"
            cursor.execute(query)
            rows = cursor.fetchall()
            column_names = [column[0] for column in cursor.description]
            df = pd.DataFrame(rows, columns=column_names)
            consolidated_df = pd.concat([consolidated_df, df], ignore_index=True)
            cursor.close()
            c.close()

        except mysql.connector.Error as error:
            raise Exception(f"Error connecting to MySQL: {error}")
        
        return consolidated_df

# ----- Page Config -----
st.set_page_config(page_title="Dashboard", page_icon=":bar_chart:", layout="wide")
st.title("Dashboard")
choice = st.sidebar.selectbox("Menu", ['Data Initialization', 'Data Analysis'])
connection = {'hostname': 'localhost', 'username': 'root', 'password': 'cdev', 'database': 'cph', 'table': 'online_survey'}

if os.path.exists('./dataset.csv'):
    df = pd.read_csv('dataset.csv', index_col=None)

if not os.path.exists("images"):
    os.mkdir("images")

if choice == 'Data Initialization':
    with st.form(key='form1'):
        hostname = st.text_input(label='Hostname:', value=connection['hostname'])
        username = st.text_input(label='Hostname:', value=connection['username'])
        password = st.text_input(label='Hostname:', value=connection['password'], type='password')
        database = st.text_input(label='Hostname:', value=connection['database'])
        table = st.text_input(label='Hostname:', value=connection['table'])
        save_config = st.form_submit_button(label='Save MySQL Connection')
        st.text('OR')
        uploaded_file = st.file_uploader("Choose a CSV file")
        save_file = st.form_submit_button(label='Save')
    
    if save_config:
        test_connection = dict()
        test_connection['hostname'] = hostname
        test_connection['username'] = username
        test_connection['password'] = password
        test_connection['database'] = database
        test_connection['table'] = table
        try:
            df = retrieve_data_from_connections(test_connection)
        except:
            st.error('Check Database details', icon="🚨")
        else:
            connection = test_connection
            st.success('Connection established!', icon='🤖')

    if save_file:
        try:
            df = pd.read_csv(uploaded_file)
            df.to_csv('dataset.csv', index=None)
        except:
            st.error('Something went wrong...', icon="🚨")
        else:
            st.success('CSV loaded!', icon='🤖')

if choice == 'Data Analysis':

    # ----- SQL Connection -----
    # r_connection = {'hostname': 'localhost', 'username': 'root', 'password': 'cdev', 'database': 'cph', 'table': 'reservations'}
    # r_df = retrieve_data_from_connections(r_connection)
    # ----- DataFrame Edits -----
    df['submission_date'] = pd.to_datetime(df['submission_date'], format='%Y-%m-%d')
    df['month'] = df['submission_date'].dt.month
    df['average_rating'] = df[['service_rating', 'room_rating', 'cleanliness_rating', 'location_rating']].mean(axis=1)

    age_bins = [0, 30, 40, 50, 60, float('inf')]  # Bins for age groups
    age_labels = ['0-30', '31-40', '41-50', '51-60', '61+']  # Labels for each age group
    df['age_group'] = pd.cut(df['age'], bins=age_bins, labels=age_labels)

    age_group_percentage = df['age_group'].value_counts(normalize=True).reset_index()
    age_group_percentage.columns = ['AgeGroup', 'Percentage']
    age_group_percentage['Percentage'] *= 100

    ihg_tier_percentage = df['IHG_tier'].value_counts(normalize=True).reset_index()
    ihg_tier_percentage.columns = ['IHG_tier', 'Percentage']
    ihg_tier_percentage['Percentage'] *= 100

    # r_df['check_in'] = pd.to_datetime(r_df['check_in'], format='%Y-%m-%d')
    # r_df['check_out'] = pd.to_datetime(r_df['check_out'], format='%Y-%m-%d')
    # r_df['date_delta'] = (r_df['check_out'] - r_df['check_in']).dt.days
    # r_df['estim_charge'] = r_df['date_delta'] * r_df['rate']

    # ----- Sidebar -----
    st.sidebar.header("Filter")
    country = st.sidebar.multiselect("Countries of origin: ", options=df['country'].unique(), default=df['country'].unique())
    ihg_tier = st.sidebar.multiselect("IHG Tier: ", options=df['IHG_tier'].unique(), default=df['IHG_tier'].unique())
    ihg_reward_recognition = st.sidebar.multiselect("IHG Reward Recognition: ", options=df['IHG_reward_recognition'].unique(), default=df['IHG_reward_recognition'].unique())

    df_selection = df.query("country == @country & IHG_tier == @ihg_tier & IHG_reward_recognition == @ihg_reward_recognition")

    # ----- KPI -----
    # total_sales = r_df['estim_charge'].sum()
    total_sales = 9000

    current_month = 4
    # current_month = datetime.now().month
    if current_month != 1:
        previous_month = current_month - 1
    else:
        previous_month = 12

    average_service_rating = df[df['month'] == current_month]['service_rating'].mean()
    average_service_delta = df[df['month'] == current_month]['service_rating'].mean() - df[df['month'] == previous_month]['service_rating'].mean()

    average_room_rating = df[df['month'] == current_month]['room_rating'].mean()
    average_room_delta = df[df['month'] == current_month]['room_rating'].mean() - df[df['month'] == previous_month]['room_rating'].mean()

    average_cleanliness_rating = df[df['month'] == current_month]['cleanliness_rating'].mean()
    average_cleanliness_delta = df[df['month'] == current_month]['cleanliness_rating'].mean() - df[df['month'] == previous_month]['cleanliness_rating'].mean()

    kpi1, kpi2, kpi3, kpi4 = st.columns(4)
    kpi1.metric(
        label="Projected Revenue 💸",
        value=total_sales,
        delta=131,
    )
    kpi2.metric(
        label="Average Service Rating This Month🛎️",
        value=average_service_rating,
        delta=round(average_service_delta, 2)
    )
    kpi3.metric(
        label="Average Room Rating This Month🚪",
        value=average_room_rating,
        delta=round(average_room_delta, 2)
    )
    kpi4.metric(
        label="Average Cleanliness Rating This Month🧼",
        value=average_cleanliness_rating,
        delta=round(average_cleanliness_delta, 2)
    )
    st.markdown("***")

    # ----- Charts -----


    rating_by_month = st.selectbox("Rating by month: ", options=['service_rating', 'room_rating', 'cleanliness_rating', 'location_rating'])
    fig_zero = px.histogram(df_selection[['service_rating', 'room_rating', 'cleanliness_rating', 'location_rating', 'IHG_tier', 'month']], x='month', y=rating_by_month, height=400, histfunc='avg')
    fig_zero.update_layout(bargap=0.2)
    fig_zero.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 1
        )
    )
    st.write(fig_zero)

    score_by_month = st.selectbox("Score by month: ", options=['amenities', 'sleep_quality'])
    fig_one = px.bar(data_frame=df_selection[['amenities', 'sleep_quality', 'month']].groupby(['month']).mean().reset_index(), x='month', y=score_by_month)
    fig_one.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 1
        ),
        yaxis_range=[0,1]
    )
    st.write(fig_one)

    age_to_rating = st.selectbox("Age to rating: ", options=['average_rating', 'service_rating', 'room_rating', 'cleanliness_rating', 'location_rating'])
    st.markdown('### Age to Rating')
    fig_two = px.line(df_selection[['age', 'average_rating', 'service_rating', 'room_rating', 'cleanliness_rating', 'location_rating']].groupby(['age']).mean().reset_index(), x='age', y=age_to_rating)
    fig_two.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 10
        )
    )
    st.write(fig_two)

    five_one, five_two = st.columns(2)
    with five_one:
        st.markdown('### Age Demographics Pie Chart (fixed)')
        fig_nin = px.pie(age_group_percentage, values='Percentage', names='AgeGroup')
        st.write(fig_nin)
    with five_two:
        st.markdown('### IHG Tier Demographic Pie Chart (fixed)')
        fig_ten = px.pie(ihg_tier_percentage, values='Percentage', names='IHG_tier')
        st.write(fig_ten)

    text_punc = " ".join(review for review in df_selection['review'])
    text = text_punc.replace('.', '')
    text = text.replace(',', '')
    text = text.replace('!', '')
    wordcloud = WordCloud(width=2000, height=1000, stopwords=['good', 'the', 'a', 'uhhh', 'was', 'it', 'on', 'my', 'for', 'Not', 'it', 'uhhhh', 'even', 'in', 'very', 'were', 'be', 'our', 'me', 'Gave', 'to', 'that', 'could', 'and', 'Needs', 'though', 'any', 'an', 'with', 'us', 'did', 'What', 'While', 'too', 'at']).generate(text)
    plt.figure(figsize=(20,10))
    plt.imshow(wordcloud, interpolation='bilinear')
    st.pyplot(plt.gcf())
    st.text('to save this figure, right click and press "save image as"')
    st.dataframe(df_selection)

    # BACKEND FIGS
    backend_fig1_1 = px.histogram(df_selection[['service_rating', 'room_rating', 'cleanliness_rating', 'location_rating', 'IHG_tier', 'month']], x='month', y='service_rating', height=400, histfunc='avg')
    backend_fig1_1.update_layout(bargap=0.2)
    backend_fig1_1.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 1
        )
    )
    backend_fig1_2 = px.histogram(df_selection[['service_rating', 'room_rating', 'cleanliness_rating', 'location_rating', 'IHG_tier', 'month']], x='month', y='room_rating', height=400, histfunc='avg')
    backend_fig1_2.update_layout(bargap=0.2)
    backend_fig1_2.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 1
        )
    )
    backend_fig1_3 = px.histogram(df_selection[['service_rating', 'room_rating', 'cleanliness_rating', 'location_rating', 'IHG_tier', 'month']], x='month', y='cleanliness_rating', height=400, histfunc='avg')
    backend_fig1_3.update_layout(bargap=0.2)
    backend_fig1_3.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 1
        )
    )
    backend_fig1_4 = px.histogram(df_selection[['service_rating', 'room_rating', 'cleanliness_rating', 'location_rating', 'IHG_tier', 'month']], x='month', y='location_rating', height=400, histfunc='avg')
    backend_fig1_4.update_layout(bargap=0.2)
    backend_fig1_4.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 1
        )
    )

    backend_fig2_1 = px.bar(data_frame=df_selection[['amenities', 'sleep_quality', 'month']].groupby(['month']).mean().reset_index(), x='month', y='amenities')
    backend_fig2_1.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 1
        ),
        yaxis_range=[0,1]
    )
    backend_fig2_2 = px.bar(data_frame=df_selection[['amenities', 'sleep_quality', 'month']].groupby(['month']).mean().reset_index(), x='month', y='sleep_quality')
    backend_fig2_2.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 1
        ),
        yaxis_range=[0,1]
    )

    backend_fig3_1 = px.line(df_selection[['age', 'average_rating', 'service_rating', 'room_rating', 'cleanliness_rating', 'location_rating']].groupby(['age']).mean().reset_index(), x='age', y='average_rating')
    backend_fig3_1.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 10
        )
    )
    backend_fig3_2 = px.line(df_selection[['age', 'average_rating', 'service_rating', 'room_rating', 'cleanliness_rating', 'location_rating']].groupby(['age']).mean().reset_index(), x='age', y='service_rating')
    backend_fig3_2.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 10
        )
    )
    backend_fig3_3 = px.line(df_selection[['age', 'average_rating', 'service_rating', 'room_rating', 'cleanliness_rating', 'location_rating']].groupby(['age']).mean().reset_index(), x='age', y='room_rating')
    backend_fig3_3.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 10
        )
    )
    backend_fig3_4 = px.line(df_selection[['age', 'average_rating', 'service_rating', 'room_rating', 'cleanliness_rating', 'location_rating']].groupby(['age']).mean().reset_index(), x='age', y='cleanliness_rating')
    backend_fig3_4.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 10
        )
    )
    backend_fig3_5 = px.line(df_selection[['age', 'average_rating', 'service_rating', 'room_rating', 'cleanliness_rating', 'location_rating']].groupby(['age']).mean().reset_index(), x='age', y='location_rating')
    backend_fig3_5.update_layout(
        xaxis = dict(
            tickmode = 'linear',
            dtick = 10
        )
    )

    figures_list = [backend_fig1_1, backend_fig1_2, backend_fig1_3, backend_fig1_4, backend_fig2_1, backend_fig2_2, backend_fig3_1, backend_fig3_2, backend_fig3_3, backend_fig3_4, backend_fig3_5]
    images_to_pptx = st.button("Save Figures into PPTX")
    output_data = None
    if images_to_pptx:
        try:
            download_plotly_figures_as_png(figures_list, "images")
            output_data = create_binaryrep_with_images("images")
        except:
            st.error('Something went wrong', icon="🚨")
        else:
            st.success('Images saved!', icon='🤖')

    if output_data:
        download_pptx = st.download_button(
            label='📥 Download',
            data = output_data,
            file_name = 'my_power.pptx'
        )